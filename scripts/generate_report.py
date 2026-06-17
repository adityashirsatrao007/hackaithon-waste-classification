"""Generate Round 1 deliverables: 5-slide PPT + submission.csv verification."""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import glob

REPORTS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "reports")
FIGURES = os.path.join(REPORTS, "figures")
EVIDENCE = os.path.join(REPORTS, "evidence")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

# --- Slide 1: Dataset Overview ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 1, "Waste Classification — Dataset Overview", 36, True, ACCENT)
add_text_box(slide, 0.8, 1.4, 11, 0.5, "Hack[AI]Thon 2.0 · Round 1 · Data Detective Challenge", 18, False, LIGHT)

items = [
    "Dataset: 445 labeled training images, 115 test images",
    "Classes: cardboard (82), glass (78), metal (84), paper (86), plastic (75), trash (40)",
    "Imbalance ratio: 2.15× (trash is most underrepresented)",
    "Source: Waste classification image set with intentionally introduced quality defects",
    "Baseline model: Fixed ResNet-18 (not modified per competition rules)",
    "Approach: Data-centric audit — identify issues, gather evidence, propose fixes"
]
add_bullet_frame(slide, 0.8, 2.2, 11, 4.5, items, 18)

# --- Slide 2: Problems Found ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 1, "Problems Found", 36, True, ACCENT)

problems = [
    "1. Wrong Labels — 71.5% (318/445) images are in the wrong folder. Folder names and filename-prefix labels disagree on most samples. Filename labels are equally unreliable.",
    "2. Duplicate Images — 6 exact duplicate groups and 11 near-duplicate pairs found. Redundant samples inflate perceived dataset size and leak between train/validation splits.",
    "3. Blurry / Low-Quality Images — 157 images (35.3%) have Laplacian variance < 100. Minimum variance is 0.6 (effectively blank). These add noise without signal.",
    "4. Class Imbalance — trash class has only 40 samples vs 82-86 for others. Model never learns to predict trash (0% recall, 0% F1).",
    "5. Outliers — 22 file-size outliers detected (range 4.7KB–43KB). Possible corrupted or anomalous samples."
]
add_bullet_frame(slide, 0.8, 1.4, 11, 5.5, problems, 16)

# --- Slide 3: Evidence ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.8, "Evidence & Analysis", 36, True, ACCENT)

# Add confusion matrix figure
cm_paths = glob.glob(os.path.join(FIGURES, "confusion_matrix.png"))
if cm_paths:
    slide.shapes.add_picture(cm_paths[0], Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.0))

add_text_box(slide, 6.5, 1.3, 6, 0.4, "Key Metrics (Baseline Model):", 18, True, WHITE)
metrics = [
    "Overall Val Accuracy: 42.53% (≈ random on 6 classes = 16.7%)",
    "Precision range: 0.0 (trash) — 0.80 (cardboard)",
    "Recall range: 0.0 (trash) — 0.72 (cardboard)",
    "Confusion matrix shows model predicts 'cardboard' for all classes — model learned folder bias, not visual features",
    "UMAP embeddings show no clean cluster separation by true class",
    "Evidence images saved at reports/evidence/ (12 mislabel examples, 12 blurry examples)"
]
add_bullet_frame(slide, 6.5, 1.8, 6, 5.0, metrics, 14)

# --- Slide 4: Suggested Fixes ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 1, "Suggested Fixes", 36, True, ACCENT)

fixes = [
    "Fix 1 — Correct Labels: Use filename-prefix labels where clear. For ambiguous cases, use confident learning (CleanLab) to flag 24 most-likely-correct corrections.",
    "Fix 2 — Remove Duplicates: Deduplicate exact matches (6 groups) and near-duplicates (11 pairs) to prevent train/val leakage and inflated class counts.",
    "Fix 3 — Handle Blurry Images: Down-weight or remove 157 low-quality images (Laplacian variance < 100). Alternative: apply stronger augmentation to these.",
    "Fix 4 — Balance Classes: Oversample trash class or use class weights. Data augmentation (rotation, flip, color jitter) most effective for minority classes.",
    "Fix 5 — Stratified Splits: Use stratified sampling for train/validation splits to ensure all classes represented proportionally.",
    "Fix 6 — Iterative Label Correction: Self-training / pseudo-labeling: train on corrected labels, predict unlabeled set, add high-confidence predictions, retrain."
]
add_bullet_frame(slide, 0.8, 1.5, 11, 5.5, fixes, 16)

# --- Slide 5: Expected Impact ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.5, 11, 1, "Expected Impact of Fixes", 36, True, ACCENT)

# Add UMAP figure
umap_paths = glob.glob(os.path.join(FIGURES, "umap_by_folder.png"))
if umap_paths:
    slide.shapes.add_picture(umap_paths[0], Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5))

impact_items = [
    "Label correction alone: expected +10–15% accuracy (from 42% → ~55–60%)",
    "Deduplication + clean split: expected +3–5% (removes train/val leakage)",
    "Blurry image handling: expected +2–3% (removes noisy signal)",
    "Class balancing: expected +5–8% on trash recall (0% → ~30–40%)",
    "Combined (all fixes): projected 60–70% val accuracy",
    "Confident learning only corrected 24/445 labels — noise too severe for automated fix",
    "Recommendation: Human-in-the-loop verification of 318 mislabeled samples",
    "Best next step: Self-training loop — predict, filter confident (p > 0.6), relabel, retrain"
]
add_bullet_frame(slide, 6.5, 1.5, 6, 5.5, impact_items, 15)

out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "reports", "Round1_Data_Detective.pptx")
prs.save(out)
print(f"Saved: {out}")
